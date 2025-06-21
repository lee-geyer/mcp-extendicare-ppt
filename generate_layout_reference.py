#!/usr/bin/env python3
"""
Generate a visual reference deck showing all Extendicare layouts with placeholder IDs and types
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
from pathlib import Path
from datetime import datetime

def generate_layout_reference(template_path, output_path="Extendicare_Layout_Reference.pptx"):
    """Generate a reference presentation showing all layouts with placeholder details"""
    
    # Load the template
    template_prs = Presentation(template_path)
    
    # Create new presentation using the template
    prs = Presentation(template_path)
    
    # Clear any existing slides
    while len(prs.slides) > 0:
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[0])
    
    print(f"Generating reference deck for {len(template_prs.slide_layouts)} layouts...")
    
    # For each layout, create a slide showing all placeholders
    for layout_idx, layout in enumerate(template_prs.slide_layouts):
        print(f"Processing Layout {layout_idx}: {layout.name}")
        
        # Add slide using this layout
        slide = prs.slides.add_slide(layout)
        
        # Fill each placeholder with its information
        for placeholder in slide.placeholders:
            try:
                ph_idx = placeholder.placeholder_format.idx
                ph_type = str(placeholder.placeholder_format.type)
                ph_name = placeholder.name
                
                # Create info text
                info_text = f"{ph_name}\n{ph_type}\nidx: {ph_idx}"
                
                # Add text based on placeholder type
                if hasattr(placeholder, 'text_frame'):
                    placeholder.text = info_text
                    
                    # Style the text to make it readable
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(14)
                            run.font.bold = True
                            # Use a color that stands out
                            run.font.color.rgb = RGBColor(0, 0, 128)  # Dark blue
                    
                    # Center align for better visibility
                    if placeholder.text_frame.paragraphs:
                        placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                elif 'PICTURE' in ph_type:
                    # For picture placeholders, we'll skip or could add a placeholder image
                    # with the text overlaid
                    pass
                
            except Exception as e:
                print(f"  - Error processing placeholder {placeholder.name}: {e}")
        
        # Add layout info to the slide (if there's a title placeholder)
        for placeholder in slide.placeholders:
            if 'TITLE' in str(placeholder.placeholder_format.type):
                current_text = placeholder.text
                placeholder.text = f"Layout {layout_idx}: {layout.name}\n{current_text}"
                break
    
    # Add an index slide at the beginning
    if len(template_prs.slide_layouts) > 0:
        # Use the first layout for the index
        first_slide = prs.slides.add_slide(template_prs.slide_layouts[0])
        
        # Find title placeholder
        for placeholder in first_slide.placeholders:
            if 'TITLE' in str(placeholder.placeholder_format.type):
                placeholder.text = "Extendicare Layout Reference Guide"
            elif 'SUBTITLE' in str(placeholder.placeholder_format.type) or 'BODY' in str(placeholder.placeholder_format.type):
                if hasattr(placeholder, 'text_frame'):
                    placeholder.text = f"Total Layouts: {len(template_prs.slide_layouts)}\nGenerated: {datetime.now().strftime('%B %d, %Y')}\n\nEach slide shows placeholder names, types, and idx values"
                break
        
        # Move the index slide to the beginning
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(0, slides[-1])
    
    # Save the reference presentation
    prs.save(output_path)
    print(f"\nReference deck saved to: {output_path}")
    print(f"Total slides generated: {len(prs.slides)}")
    
    # Also create a summary JSON with layout info
    layout_summary = []
    for idx, layout in enumerate(template_prs.slide_layouts):
        layout_info = {
            "index": idx,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders),
            "placeholders": []
        }
        
        for ph in layout.placeholders:
            ph_info = {
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
                "has_text_frame": hasattr(ph, 'text_frame')
            }
            layout_info["placeholders"].append(ph_info)
        
        layout_summary.append(layout_info)
    
    # Save summary
    summary_path = "layout_reference_summary.json"
    with open(summary_path, "w") as f:
        json.dump(layout_summary, f, indent=2)
    print(f"Layout summary saved to: {summary_path}")

if __name__ == "__main__":
    # Use the actual template path
    template_path = "/Users/leegeyer/Library/CloudStorage/OneDrive-Extendicare(Canada)Inc/Investor Quarterly Reporting/ExtendicareTemplate.pptx"
    
    if not Path(template_path).exists():
        print(f"Error: Template not found at {template_path}")
        print("Please update the template_path variable with the correct path")
    else:
        generate_layout_reference(template_path)
        print("\nYou can now open the generated 'Extendicare_Layout_Reference.pptx' file")
        print("to see all layouts with their placeholder IDs and types labeled!")