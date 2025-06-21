from pptx import Presentation
import json
from pathlib import Path

def analyze_template(template_path):
    """Analyze PowerPoint template and extract layout information"""
    
    # Load the template (use .pptx version)
    prs = Presentation(template_path)
    
    # Get basic info
    print(f"Total slide layouts: {len(prs.slide_layouts)}")
    print("=" * 80)
    
    # List all layouts with their names and placeholder info
    layouts_info = []
    
    for idx, layout in enumerate(prs.slide_layouts):
        layout_data = {
            "index": idx,
            "name": layout.name,
            "placeholders": []
        }
        
        for placeholder in layout.placeholders:
            ph_info = {
                "idx": placeholder.placeholder_format.idx,
                "type": str(placeholder.placeholder_format.type),
                "name": placeholder.name,
                "has_text_frame": hasattr(placeholder, 'text_frame')
            }
            layout_data["placeholders"].append(ph_info)
        
        layouts_info.append(layout_data)
        
        # Print layout info
        print(f"\nLayout {idx}: {layout.name}")
        print(f"  Placeholders: {len(layout.placeholders)}")
        for ph in layout_data["placeholders"]:
            print(f"    - idx: {ph['idx']}, name: '{ph['name']}', type: {ph['type']}")
    
    # Save to JSON
    output_file = "template_analysis.json"
    with open(output_file, "w") as f:
        json.dump(layouts_info, f, indent=2)
    
    print(f"\n{'=' * 80}")
    print(f"Analysis saved to {output_file}")
    
    return layouts_info

if __name__ == "__main__":
    # UPDATE THIS PATH to your .pptx template
    template_path = "../ExtendicareTemplate.pptx"
    
    if not Path(template_path).exists():
        print(f"Error: Template file not found at {template_path}")
        print("Please update the template_path variable")
    else:
        analyze_template(template_path)