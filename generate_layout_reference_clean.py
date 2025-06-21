#!/usr/bin/env python3
"""
Generate a clean layout reference by creating a new presentation with template layouts
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
from pathlib import Path
from datetime import datetime
from PIL import Image, ImageDraw, ImageFont
import io
import shutil

def create_placeholder_image(info_text, width=400, height=300):
    """Create a placeholder image with the reference info"""
    # Create image with light blue background
    img = Image.new('RGB', (width, height), color='#E8F4F8')
    draw = ImageDraw.Draw(img)
    
    # Add border
    draw.rectangle([(0, 0), (width-1, height-1)], outline='#0066CC', width=3)
    
    # Add text
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
    except:
        font = ImageFont.load_default()
    
    # Split text into lines
    lines = info_text.split('\n')
    y_offset = 20
    
    for line in lines:
        text_bbox = draw.textbbox((0, 0), line, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        x = (width - text_width) // 2
        draw.text((x, y_offset), line, fill='#003366', font=font)
        y_offset += 30
    
    # Save to BytesIO
    img_io = io.BytesIO()
    img.save(img_io, 'PNG')
    img_io.seek(0)
    
    return img_io

def generate_layout_reference(template_path, output_path="Extendicare_Layout_Reference_Clean.pptx"):
    """Generate a reference presentation using template layouts in a clean way"""
    
    print(f"Loading template from: {template_path}")
    
    # First, let's create a working copy of the template
    working_template = "working_template.pptx"
    shutil.copy(template_path, working_template)
    
    # Load the working template
    template_prs = Presentation(working_template)
    
    # Count layouts
    num_layouts = len(template_prs.slide_layouts)
    print(f"Found {num_layouts} layouts in template")
    
    # Create a completely new, empty presentation
    prs = Presentation()
    
    # We'll manually copy key properties from the template
    prs.slide_width = template_prs.slide_width
    prs.slide_height = template_prs.slide_height
    
    print("\nCreating manual reference slides...")
    
    # Since we can't reliably use the template layouts directly,
    # let's create informative slides about each layout
    for layout_idx, layout in enumerate(template_prs.slide_layouts):
        print(f"Processing Layout {layout_idx}: {layout.name}")
        
        # Use the blank layout from the new presentation
        blank_layout = prs.slide_layouts[5]  # Usually the blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Add title
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = Inches(9)
        title_height = Inches(1)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = f"Layout {layout_idx}: {layout.name}"
        
        # Format title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(204, 0, 0)
        
        # Create a visual representation of placeholders
        y_position = Inches(1.5)
        
        for placeholder in layout.placeholders:
            ph_idx = placeholder.placeholder_format.idx
            ph_type = str(placeholder.placeholder_format.type)
            ph_name = placeholder.name
            
            # Create info box for this placeholder
            info_left = Inches(0.5)
            info_width = Inches(8.5)
            info_height = Inches(0.5)
            
            info_box = slide.shapes.add_textbox(info_left, y_position, info_width, info_height)
            info_frame = info_box.text_frame
            
            # Format the placeholder info
            info_text = f"• {ph_name} | Type: {ph_type} | idx: {ph_idx}"
            info_frame.text = info_text
            
            # Style the text
            for paragraph in info_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = 'Consolas'
                
                # Color code by type
                if 'TITLE' in ph_type:
                    paragraph.font.color.rgb = RGBColor(128, 0, 0)
                elif 'PICTURE' in ph_type or 'OBJECT' in ph_type:
                    paragraph.font.color.rgb = RGBColor(0, 128, 0)
                elif 'CHART' in ph_type:
                    paragraph.font.color.rgb = RGBColor(0, 0, 128)
                else:
                    paragraph.font.color.rgb = RGBColor(64, 64, 64)
            
            y_position += Inches(0.4)
            
            # Prevent overflow to next slide
            if y_position > Inches(6.5):
                # Add continuation text
                cont_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(8), Inches(0.5))
                cont_box.text_frame.text = "... (more placeholders exist)"
                break
    
    # Add an index slide at the beginning
    print("Adding index slide...")
    blank_layout = prs.slide_layouts[5]
    index_slide = prs.slides.add_slide(blank_layout)
    
    # Add title to index
    title_box = index_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Extendicare Layout Reference Guide"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add description
    desc_box = index_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(8.5), Inches(4))
    desc_frame = desc_box.text_frame
    desc_frame.text = (f"Reference guide for {num_layouts} Extendicare PowerPoint layouts\n\n"
                      f"Each slide shows:\n"
                      f"• Layout name and index\n"
                      f"• All placeholders with their:\n"
                      f"  - Name\n"
                      f"  - PowerPoint type code\n"
                      f"  - Index (idx) value\n\n"
                      f"Color coding:\n"
                      f"• Red = Title placeholders\n"
                      f"• Green = Picture/Object placeholders\n"
                      f"• Blue = Chart placeholders\n"
                      f"• Gray = Other placeholders\n\n"
                      f"Generated: {datetime.now().strftime('%B %d, %Y')}")
    
    for paragraph in desc_frame.paragraphs:
        paragraph.font.size = Pt(16)
    
    # Move index to beginning
    xml_slides = prs.slides._sldIdLst
    index_id = xml_slides[-1]
    xml_slides.remove(index_id)
    xml_slides.insert(0, index_id)
    
    # Clean up working template
    Path(working_template).unlink()
    
    # Save the clean presentation
    prs.save(output_path)
    print(f"\nClean reference deck saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    
    # Also save as a detailed text report
    report_path = output_path.replace('.pptx', '_report.txt')
    with open(report_path, 'w') as f:
        f.write("EXTENDICARE LAYOUT REFERENCE\n")
        f.write("=" * 60 + "\n\n")
        
        for layout_idx, layout in enumerate(template_prs.slide_layouts):
            f.write(f"LAYOUT {layout_idx}: {layout.name}\n")
            f.write("-" * 60 + "\n")
            
            for placeholder in layout.placeholders:
                ph_idx = placeholder.placeholder_format.idx
                ph_type = str(placeholder.placeholder_format.type)
                ph_name = placeholder.name
                f.write(f"  {ph_name:<30} Type: {ph_type:<20} idx: {ph_idx}\n")
            
            f.write("\n")
    
    print(f"Detailed report saved to: {report_path}")

if __name__ == "__main__":
    # Use the actual template path
    template_path = "/Users/leegeyer/Library/CloudStorage/OneDrive-Extendicare(Canada)Inc/Investor Quarterly Reporting/ExtendicareTemplate.pptx"
    
    if not Path(template_path).exists():
        print(f"Error: Template not found at {template_path}")
        print("Please update the template_path variable with the correct path")
    else:
        generate_layout_reference(template_path)