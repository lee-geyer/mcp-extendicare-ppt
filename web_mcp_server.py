#!/usr/bin/env python3
"""
Web-based MCP server for Claude Desktop integration
"""
import asyncio
import json
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime

from mcp.server import Server
from mcp.types import Tool
from pydantic import BaseModel, Field
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io
import uvicorn
from starlette.applications import Starlette
from starlette.responses import JSONResponse
from starlette.routing import Route
import sys
import os

# Add current directory to path
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

# Configuration
class Config(BaseModel):
    template_path: str = Field(description="Path to Extendicare PowerPoint template")
    output_dir: str = Field(default=".", description="Output directory for generated presentations")

# Load template analysis
def load_template_analysis():
    with open("template_analysis.json", "r") as f:
        return json.load(f)

class ExtendicarePPTServer:
    def __init__(self, config: Config):
        self.config = config
        self.template_path = Path(config.template_path)
        self.output_dir = Path(config.output_dir)
        self.layouts = load_template_analysis()
        
        # Verify template exists
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {self.template_path}")
    
    def create_placeholder_image(self, description: str, width: int = 800, height: int = 600) -> io.BytesIO:
        """Create a placeholder image with description text"""
        img = Image.new('RGB', (width, height), color='#E0E0E0')
        draw = ImageDraw.Draw(img)
        
        # Add border
        border_width = 2
        draw.rectangle(
            [(0, 0), (width-1, height-1)],
            outline='#808080',
            width=border_width
        )
        
        # Add text
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
        except:
            font = ImageFont.load_default()
        
        text_bbox = draw.textbbox((0, 0), description, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        text_height = text_bbox[3] - text_bbox[1]
        x = (width - text_width) // 2
        y = (height - text_height) // 2
        
        draw.text((x, y), description, fill='#404040', font=font)
        
        img_io = io.BytesIO()
        img.save(img_io, 'PNG')
        img_io.seek(0)
        
        return img_io
    
    async def query_layouts(self, features: Optional[List[str]] = None) -> List[Dict[str, Any]]:
        """Query layouts by features"""
        results = []
        
        for layout in self.layouts:
            if features:
                layout_name_lower = layout['name'].lower()
                placeholder_types = [p['type'].lower() for p in layout['placeholders']]
                
                matches = False
                for feature in features:
                    feature_lower = feature.lower()
                    if (feature_lower in layout_name_lower or
                        any(feature_lower in ptype for ptype in placeholder_types)):
                        matches = True
                        break
                
                if matches:
                    results.append(layout)
            else:
                results.append(layout)
        
        return results
    
    async def create_presentation(self, spec: Dict[str, Any]) -> str:
        """Create a complete presentation based on specification"""
        prs = Presentation(str(self.template_path))
        
        # Simple implementation - create one slide for demo
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Fill title
        title_placeholder = slide.shapes.title
        if title_placeholder:
            title_placeholder.text = spec.get('title', 'Generated Presentation')
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"Extendicare_Presentation_{timestamp}.pptx"
        output_path = self.output_dir / filename
        prs.save(str(output_path))
        
        return str(output_path)

# Initialize server
config = Config(
    template_path="../ExtendicareTemplate.pptx",
    output_dir="."
)

try:
    ppt_server = ExtendicarePPTServer(config)
    print("✅ Extendicare PPT Server initialized successfully")
except Exception as e:
    print(f"❌ Failed to initialize server: {e}")
    ppt_server = None

# Web endpoints
async def health_check(request):
    return JSONResponse({"status": "healthy", "service": "Extendicare PPT MCP Server"})

async def list_tools(request):
    tools = [
        {
            "name": "query_layouts",
            "description": "Search for PowerPoint layouts by features",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "features": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Features to search for (e.g., 'chart', 'picture', '2 column')"
                    }
                }
            }
        },
        {
            "name": "create_presentation",
            "description": "Create a PowerPoint presentation with Extendicare branding",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "Presentation title"
                    },
                    "content": {
                        "type": "string", 
                        "description": "Presentation content"
                    }
                }
            }
        }
    ]
    return JSONResponse({"tools": tools})

async def call_tool(request):
    try:
        data = await request.json()
        tool_name = data.get("name")
        arguments = data.get("arguments", {})
        
        if not ppt_server:
            return JSONResponse({"error": "Server not initialized"}, status_code=500)
        
        if tool_name == "query_layouts":
            features = arguments.get("features")
            results = await ppt_server.query_layouts(features)
            return JSONResponse({"result": results})
        
        elif tool_name == "create_presentation":
            output_path = await ppt_server.create_presentation(arguments)
            return JSONResponse({"result": {"success": True, "path": output_path}})
        
        else:
            return JSONResponse({"error": f"Unknown tool: {tool_name}"}, status_code=400)
    
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

# Create Starlette app
app = Starlette(
    routes=[
        Route("/", health_check),
        Route("/health", health_check),
        Route("/tools", list_tools),
        Route("/call", call_tool, methods=["POST"]),
    ]
)

if __name__ == "__main__":
    print("🚀 Starting Extendicare PPT MCP Web Server")
    print("📍 Server will be available at: http://localhost:8001")
    print("🔧 Health check: http://localhost:8001/health")
    print("🛠️  Tools endpoint: http://localhost:8001/tools")
    
    uvicorn.run(app, host="0.0.0.0", port=8001, log_level="info")