import asyncio
import json
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime

from mcp.server import Server
from mcp.types import (
    CallToolResult,
    Tool
)
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
import io

# Configuration
class Config(BaseModel):
    template_path: str = Field(description="Path to Extendicare PowerPoint template")
    output_dir: str = Field(default=".", description="Output directory for generated presentations")

# Load template analysis
def load_template_analysis():
    with open("template_analysis.json", "r") as f:
        return json.load(f)

class ExtendicarePPTServer:
    def __init__(self, config: Config):
        self.config = config
        self.template_path = Path(config.template_path)
        self.output_dir = Path(config.output_dir)
        self.layouts = load_template_analysis()
        
        # Verify template exists
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {self.template_path}")
    
    def create_placeholder_image(self, description: str, width: int = 800, height: int = 600) -> io.BytesIO:
        """Create a placeholder image with description text"""
        # Create image with light gray background
        img = Image.new('RGB', (width, height), color='#E0E0E0')
        draw = ImageDraw.Draw(img)
        
        # Add border
        border_width = 2
        draw.rectangle(
            [(0, 0), (width-1, height-1)],
            outline='#808080',
            width=border_width
        )
        
        # Add text
        try:
            # Try to use a nice font, fallback to default if not available
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
        except:
            font = ImageFont.load_default()
        
        # Calculate text position (centered)
        text_bbox = draw.textbbox((0, 0), description, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        text_height = text_bbox[3] - text_bbox[1]
        x = (width - text_width) // 2
        y = (height - text_height) // 2
        
        # Draw text
        draw.text((x, y), description, fill='#404040', font=font)
        
        # Save to BytesIO
        img_io = io.BytesIO()
        img.save(img_io, 'PNG')
        img_io.seek(0)
        
        return img_io
    
    async def query_layouts(self, features: Optional[List[str]] = None) -> List[Dict[str, Any]]:
        """Query layouts by features (e.g., 'chart', 'picture', '2 column')"""
        results = []
        
        for layout in self.layouts:
            if features:
                # Check if layout matches any of the requested features
                layout_name_lower = layout['name'].lower()
                placeholder_types = [p['type'].lower() for p in layout['placeholders']]
                
                matches = False
                for feature in features:
                    feature_lower = feature.lower()
                    if (feature_lower in layout_name_lower or
                        any(feature_lower in ptype for ptype in placeholder_types)):
                        matches = True
                        break
                
                if matches:
                    results.append(layout)
            else:
                results.append(layout)
        
        return results
    
    async def get_layout_details(self, layout_index: int) -> Dict[str, Any]:
        """Get detailed information about a specific layout"""
        if 0 <= layout_index < len(self.layouts):
            return self.layouts[layout_index]
        else:
            raise ValueError(f"Layout index {layout_index} out of range (0-{len(self.layouts)-1})")
    
    async def create_slide(self, layout_index: int, content: Dict[str, Any]) -> Dict[str, Any]:
        """Create a single slide with specified content"""
        # This will be called by create_presentation
        # Returns slide data for inclusion in presentation
        return {
            "layout_index": layout_index,
            "content": content
        }
    
    async def create_presentation(self, spec: Dict[str, Any]) -> str:
        """Create a complete presentation based on specification"""
        prs = Presentation(str(self.template_path))
        
        # Parse specification - supports multiple input formats
        slides_data = self._parse_presentation_spec(spec)
        
        # Create slides
        for slide_data in slides_data:
            layout_index = slide_data.get('layout_index', 0)
            content = slide_data.get('content', {})
            
            # Get layout
            if layout_index >= len(prs.slide_layouts):
                layout_index = 0  # Fallback to first layout
            
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            
            # Fill placeholders
            for placeholder in slide.placeholders:
                ph_name = placeholder.name.lower()
                
                # Match content to placeholders
                if 'title' in ph_name and 'title' in content:
                    placeholder.text = content['title']
                elif 'subtitle' in ph_name and 'subtitle' in content:
                    placeholder.text = content['subtitle']
                elif 'date' in ph_name and 'date' in content:
                    placeholder.text = content['date']
                elif 'heading' in ph_name and 'headings' in content:
                    # Match heading by index
                    heading_num = self._extract_number(ph_name)
                    headings = content['headings']
                    if isinstance(headings, list) and heading_num < len(headings):
                        placeholder.text = headings[heading_num]
                elif 'content' in ph_name or 'text' in ph_name:
                    if 'bullets' in content:
                        self._add_bullets(placeholder, content['bullets'])
                    elif 'content' in content:
                        if isinstance(content['content'], list):
                            self._add_bullets(placeholder, content['content'])
                        else:
                            placeholder.text = str(content['content'])
                elif 'picture' in ph_name.lower() and 'image' in content:
                    # Add placeholder image
                    img_desc = content['image']
                    img_io = self.create_placeholder_image(img_desc)
                    placeholder.insert_picture(img_io)
                elif 'quote' in ph_name and 'quote' in content:
                    placeholder.text = content['quote']
                elif 'name' in ph_name and 'attribution' in content:
                    placeholder.text = content['attribution']
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"Extendicare_Presentation_{timestamp}.pptx"
        output_path = self.output_dir / filename
        prs.save(str(output_path))
        
        return str(output_path)
    
    def _parse_presentation_spec(self, spec: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Parse different input formats into slide data"""
        # Handle different input formats
        if 'slides' in spec:
            # Explicit slides array
            return spec['slides']
        elif 'markdown' in spec:
            # Parse markdown format
            return self._parse_markdown(spec['markdown'])
        elif 'description' in spec:
            # Natural language - would need NLP, for now just create basic slides
            return self._parse_natural_language(spec['description'])
        else:
            # Assume spec is a single slide
            return [spec]
    
    def _parse_markdown(self, markdown: str) -> List[Dict[str, Any]]:
        """Parse markdown into slides"""
        slides = []
        current_slide = None
        
        lines = markdown.strip().split('\n')
        for line in lines:
            if line.startswith('# '):
                # New slide with title
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    'layout_index': 0,  # Title slide
                    'content': {'title': line[2:].strip()}
                }
            elif line.startswith('## '):
                # Subtitle or section header
                if current_slide:
                    current_slide['content']['subtitle'] = line[3:].strip()
                else:
                    current_slide = {
                        'layout_index': 5,  # 1 Column
                        'content': {'title': line[3:].strip()}
                    }
            elif line.startswith('- '):
                # Bullet point
                if not current_slide:
                    current_slide = {'layout_index': 5, 'content': {}}
                if 'bullets' not in current_slide['content']:
                    current_slide['content']['bullets'] = []
                current_slide['content']['bullets'].append(line[2:].strip())
        
        if current_slide:
            slides.append(current_slide)
        
        return slides
    
    def _parse_natural_language(self, description: str) -> List[Dict[str, Any]]:
        """Basic parsing of natural language description"""
        # This is a simplified version - in production, you'd use NLP
        slides = []
        
        # Create a title slide
        slides.append({
            'layout_index': 0,
            'content': {
                'title': 'Presentation',
                'subtitle': 'Generated from description',
                'date': datetime.now().strftime("%B %d, %Y")
            }
        })
        
        # Add a content slide
        slides.append({
            'layout_index': 5,
            'content': {
                'title': 'Content',
                'content': description
            }
        })
        
        return slides
    
    def _add_bullets(self, placeholder, bullets):
        """Add bullet points to a placeholder"""
        if hasattr(placeholder, 'text_frame'):
            text_frame = placeholder.text_frame
            text_frame.clear()  # Clear existing text
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
    
    def _extract_number(self, text: str) -> int:
        """Extract number from text (e.g., 'Heading 2' -> 1)"""
        import re
        numbers = re.findall(r'\d+', text)
        return int(numbers[0]) - 1 if numbers else 0

# MCP Server setup
server = Server("mcp-extendicare-ppt")

@server.list_tools()
async def list_tools():
    return [
        Tool(
            name="query_layouts",
            description="Search for PowerPoint layouts by features",
            inputSchema={
                "type": "object",
                "properties": {
                    "features": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Features to search for (e.g., 'chart', 'picture', '2 column')"
                    }
                }
            }
        ),
        Tool(
            name="get_layout_details",
            description="Get detailed information about a specific layout",
            inputSchema={
                "type": "object",
                "properties": {
                    "layout_index": {
                        "type": "integer",
                        "description": "Index of the layout (0-21)"
                    }
                },
                "required": ["layout_index"]
            }
        ),
        Tool(
            name="create_presentation",
            description="Create a PowerPoint presentation with Extendicare branding",
            inputSchema={
                "type": "object",
                "properties": {
                    "spec": {
                        "type": "object",
                        "description": "Presentation specification (supports multiple formats)"
                    }
                },
                "required": ["spec"]
            }
        )
    ]

# Global server instance
ppt_server = None

@server.call_tool()
async def call_tool(name: str, arguments: Any):
    global ppt_server
    
    try:
        if name == "query_layouts":
            features = arguments.get("features")
            results = await ppt_server.query_layouts(features)
            return results
        
        elif name == "get_layout_details":
            layout_index = arguments["layout_index"]
            details = await ppt_server.get_layout_details(layout_index)
            return details
        
        elif name == "create_presentation":
            spec = arguments["spec"]
            output_path = await ppt_server.create_presentation(spec)
            return {"success": True, "path": output_path}
        
        else:
            raise ValueError(f"Unknown tool: {name}")
    
    except Exception as e:
        raise RuntimeError(f"Error executing {name}: {str(e)}")

async def main():
    from mcp.server.stdio import stdio_server
    global ppt_server
    
    # Load config
    config = Config(
        template_path="../ExtendicareTemplate.pptx",
        output_dir="."
    )
    
    ppt_server = ExtendicarePPTServer(config)
    
    # Run the server via stdio
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            server.create_initialization_options()
        )

if __name__ == "__main__":
    asyncio.run(main())