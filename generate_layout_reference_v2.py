#!/usr/bin/env python3
"""
Generate a visual reference deck showing all Extendicare layouts with placeholder IDs and types
Version 2: Creates a clean presentation without template content
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
from pathlib import Path
from datetime import datetime

def generate_layout_reference(template_path, output_path="Extendicare_Layout_Reference_v2.pptx"):
    """Generate a reference presentation showing all layouts with placeholder details"""
    
    # First, analyze the template
    template_prs = Presentation(template_path)
    print(f"Analyzing template with {len(template_prs.slide_layouts)} layouts...")
    
    # Create a completely new presentation
    # We'll start fresh and just reference the template for layout info
    reference_prs = Presentation()
    
    # Get the default slide width/height
    reference_prs.slide_width = template_prs.slide_width
    reference_prs.slide_height = template_prs.slide_height
    
    # Process each layout from the template
    layout_data = []
    for layout_idx, layout in enumerate(template_prs.slide_layouts):
        print(f"Processing Layout {layout_idx}: {layout.name}")
        
        # Collect placeholder info for this layout
        layout_info = {
            "index": layout_idx,
            "name": layout.name,
            "placeholders": []
        }
        
        # Create a text representation of the layout
        layout_text = f"LAYOUT {layout_idx}: {layout.name}\n"
        layout_text += "=" * 50 + "\n\n"
        
        for placeholder in layout.placeholders:
            ph_idx = placeholder.placeholder_format.idx
            ph_type = str(placeholder.placeholder_format.type)
            ph_name = placeholder.name
            
            # Add to our data
            ph_info = {
                "idx": ph_idx,
                "type": ph_type,
                "name": ph_name,
                "has_text_frame": hasattr(placeholder, 'text_frame')
            }
            layout_info["placeholders"].append(ph_info)
            
            # Add to text representation
            layout_text += f"Placeholder: {ph_name}\n"
            layout_text += f"  Type: {ph_type}\n"
            layout_text += f"  Index: {ph_idx}\n"
            layout_text += "-" * 30 + "\n"
        
        layout_data.append(layout_info)
        
        # Since we can't reliably add slides with the template layouts,
        # let's create a simple text slide for each layout
        slide_layout = reference_prs.slide_layouts[5]  # Use the blank layout
        slide = reference_prs.slides.add_slide(slide_layout)
        
        # Add a text box with the layout info
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(6.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = layout_text
        
        # Format the text
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.name = 'Courier New'  # Monospace for better alignment
    
    # Save the reference presentation
    reference_prs.save(output_path)
    print(f"\nReference deck saved to: {output_path}")
    print(f"Total slides generated: {len(reference_prs.slides)}")
    
    # Also save the layout data as JSON
    summary_path = "layout_reference_summary_v2.json"
    with open(summary_path, "w") as f:
        json.dump(layout_data, f, indent=2)
    print(f"Layout summary saved to: {summary_path}")
    
    # Create a more detailed text report
    report_path = "layout_reference_report.txt"
    with open(report_path, "w") as f:
        f.write("EXTENDICARE POWERPOINT TEMPLATE LAYOUT REFERENCE\n")
        f.write("=" * 60 + "\n")
        f.write(f"Generated: {datetime.now().strftime('%B %d, %Y %H:%M:%S')}\n")
        f.write(f"Total Layouts: {len(layout_data)}\n\n")
        
        for layout in layout_data:
            f.write(f"\nLAYOUT {layout['index']}: {layout['name']}\n")
            f.write("-" * 60 + "\n")
            
            # Group placeholders by type
            by_type = {}
            for ph in layout['placeholders']:
                ph_type = ph['type'].split('(')[0].strip()
                if ph_type not in by_type:
                    by_type[ph_type] = []
                by_type[ph_type].append(ph)
            
            for ph_type, placeholders in sorted(by_type.items()):
                f.write(f"\n  {ph_type} Placeholders:\n")
                for ph in placeholders:
                    f.write(f"    - {ph['name']} (idx: {ph['idx']})\n")
            
            f.write("\n")
    
    print(f"Detailed report saved to: {report_path}")
    print("\nNOTE: Since the template has existing content, this version creates")
    print("text-based reference slides to avoid conflicts.")
    print("Open the PowerPoint file to see placeholder details for each layout.")

if __name__ == "__main__":
    # Use the actual template path
    template_path = "/Users/leegeyer/Library/CloudStorage/OneDrive-Extendicare(Canada)Inc/Investor Quarterly Reporting/ExtendicareTemplate.pptx"
    
    if not Path(template_path).exists():
        print(f"Error: Template not found at {template_path}")
        print("Please update the template_path variable with the correct path")
    else:
        generate_layout_reference(template_path)