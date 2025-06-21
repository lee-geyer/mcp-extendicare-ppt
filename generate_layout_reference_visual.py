#!/usr/bin/env python3
"""
Generate a visual reference deck using actual Extendicare layouts with placeholder info filled in
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
from pathlib import Path
from datetime import datetime
from PIL import Image, ImageDraw, ImageFont
import io

def create_placeholder_image(info_text, width=400, height=300):
    """Create a placeholder image with the reference info"""
    # Create image with light blue background
    img = Image.new('RGB', (width, height), color='#E8F4F8')
    draw = ImageDraw.Draw(img)
    
    # Add border
    draw.rectangle([(0, 0), (width-1, height-1)], outline='#0066CC', width=3)
    
    # Add text
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
    except:
        font = ImageFont.load_default()
    
    # Split text into lines
    lines = info_text.split('\n')
    y_offset = 20
    
    for line in lines:
        text_bbox = draw.textbbox((0, 0), line, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        x = (width - text_width) // 2
        draw.text((x, y_offset), line, fill='#003366', font=font)
        y_offset += 30
    
    # Save to BytesIO
    img_io = io.BytesIO()
    img.save(img_io, 'PNG')
    img_io.seek(0)
    
    return img_io

def generate_layout_reference(template_path, output_path="Extendicare_Layout_Reference_Visual.pptx"):
    """Generate a reference presentation using actual layouts with placeholder info"""
    
    # Load the template
    print(f"Loading template from: {template_path}")
    prs = Presentation(template_path)
    
    # Remove any existing slides (template content)
    print("Removing existing template slides...")
    xml_slides = prs.slides._sldIdLst
    slides_to_remove = list(xml_slides)
    for slide_id in slides_to_remove:
        xml_slides.remove(slide_id)
    
    print(f"Creating reference slides for {len(prs.slide_layouts)} layouts...")
    
    # For each layout, create a slide and fill placeholders with their info
    for layout_idx, layout in enumerate(prs.slide_layouts):
        print(f"Creating slide for Layout {layout_idx}: {layout.name}")
        
        # Add a new slide using this layout
        slide = prs.slides.add_slide(layout)
        
        # Fill each placeholder with its reference information
        for placeholder in slide.placeholders:
            ph_idx = placeholder.placeholder_format.idx
            ph_type = str(placeholder.placeholder_format.type)
            ph_name = placeholder.name
            
            # Create the reference text
            ref_text = f"{ph_name}\n{ph_type}\nidx: {ph_idx}"
            
            # Handle different placeholder types
            if hasattr(placeholder, 'text_frame'):
                # Text placeholder - fill with reference info
                placeholder.text = ref_text
                
                # Style the text
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
                    
                    # Center align for better visibility
                    paragraph.alignment = PP_ALIGN.CENTER
            
            elif 'PICTURE' in ph_type or 'OBJECT' in ph_type or 'CHART' in ph_type:
                # For image/object placeholders, insert a placeholder image with the info
                try:
                    # Get approximate dimensions (default if not available)
                    width = 400
                    height = 300
                    
                    # Create image with the reference info
                    img_io = create_placeholder_image(ref_text, width, height)
                    
                    # Insert the image
                    placeholder.insert_picture(img_io)
                except Exception as e:
                    print(f"  - Could not insert image for {ph_name}: {e}")
                    # Try to add text if image fails
                    if hasattr(placeholder, 'text_frame'):
                        placeholder.text = ref_text
            
            elif 'TABLE' in ph_type:
                # For table placeholders, try to add text
                if hasattr(placeholder, 'text_frame'):
                    placeholder.text = ref_text
        
        # Add layout info to the title if there's a title placeholder
        for placeholder in slide.placeholders:
            if 'TITLE' in str(placeholder.placeholder_format.type) and hasattr(placeholder, 'text_frame'):
                # Prepend layout info to existing text
                existing_text = placeholder.text
                placeholder.text = f"Layout {layout_idx}: {layout.name}\n{existing_text}"
                
                # Make the first line (layout info) stand out
                if placeholder.text_frame.paragraphs:
                    first_paragraph = placeholder.text_frame.paragraphs[0]
                    for run in first_paragraph.runs:
                        run.font.size = Pt(18)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(204, 0, 0)  # Red for layout info
                break
    
    # Add an index slide at the beginning
    if len(prs.slide_layouts) > 0:
        print("Adding index slide...")
        index_slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Fill the index slide
        for placeholder in index_slide.placeholders:
            if 'TITLE' in str(placeholder.placeholder_format.type):
                placeholder.text = "Extendicare Layout Reference Guide"
                if hasattr(placeholder, 'text_frame'):
                    for run in placeholder.text_frame.paragraphs[0].runs:
                        run.font.size = Pt(28)
                        run.font.bold = True
            elif ('SUBTITLE' in str(placeholder.placeholder_format.type) or 
                  'BODY' in str(placeholder.placeholder_format.type)):
                if hasattr(placeholder, 'text_frame'):
                    placeholder.text = (f"Visual reference showing all {len(prs.slide_layouts)} layouts\n"
                                      f"Each placeholder displays:\n"
                                      f"• Placeholder Name\n"
                                      f"• Type (with PowerPoint type code)\n"
                                      f"• Index (idx) value\n\n"
                                      f"Generated: {datetime.now().strftime('%B %d, %Y')}")
                    # Format the text
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(16)
                break
        
        # Move index slide to the beginning
        slide_id_list = prs.slides._sldIdLst
        index_slide_id = slide_id_list[-1]
        slide_id_list.remove(index_slide_id)
        slide_id_list.insert(0, index_slide_id)
    
    # Save the presentation
    prs.save(output_path)
    print(f"\nVisual reference deck saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("\nEach slide shows the actual Extendicare layout with placeholder reference info!")

if __name__ == "__main__":
    # Use the actual template path
    template_path = "/Users/leegeyer/Library/CloudStorage/OneDrive-Extendicare(Canada)Inc/Investor Quarterly Reporting/ExtendicareTemplate.pptx"
    
    if not Path(template_path).exists():
        print(f"Error: Template not found at {template_path}")
        print("Please update the template_path variable with the correct path")
    else:
        generate_layout_reference(template_path)